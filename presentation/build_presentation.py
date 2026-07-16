#!/usr/bin/env python3
"""Builds the Victor framework academic presentation (PowerPoint).

Usage:  python3 build_presentation.py
Output: Victor-App-Framework.pptx (same directory)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ------------------------------------------------------------------ palette
INK      = RGBColor(0x16, 0x21, 0x33)   # deep navy — primary text / dark bg
INK_SOFT = RGBColor(0x3B, 0x47, 0x5E)   # slate — secondary text
ACCENT   = RGBColor(0x1D, 0x4E, 0xD8)   # blue — accent
ACCENT2  = RGBColor(0x0F, 0x76, 0x6E)   # teal — secondary accent
AMBER    = RGBColor(0xB4, 0x5A, 0x09)   # amber — highlights
PAPER    = RGBColor(0xFF, 0xFF, 0xFF)
MIST     = RGBColor(0xEE, 0xF2, 0xF7)   # light panel fill
LINE     = RGBColor(0xC9, 0xD3, 0xE0)
CODE_BG  = RGBColor(0x1E, 0x28, 0x3A)
CODE_FG  = RGBColor(0xD8, 0xE2, 0xF0)
CODE_HL  = RGBColor(0x93, 0xC5, 0xFD)

SERIF = "Georgia"
SANS  = "Calibri"
MONO  = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

slide_no = 0


# ------------------------------------------------------------------ helpers
def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _shadow_off(shape):
    el = shape._element.spPr
    existing = el.find(qn('a:effectLst'))
    if existing is None:
        el.append(el.makeelement(qn('a:effectLst'), {}))


def rect(slide, x, y, w, h, color, line_color=None, line_w=None, round_=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(kind, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = line_w or Pt(0.75)
    _shadow_off(sp)
    if round_:
        try:
            sp.adjustments[0] = 0.08
        except Exception:
            pass
    return sp


def textbox(slide, x, y, w, h, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def put(tf, text, size=16, color=INK, bold=False, italic=False, font=SANS,
        align=PP_ALIGN.LEFT, space_after=6, first=False, level=0,
        line_spacing=1.0):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.level = level
    p.line_spacing = line_spacing
    if isinstance(text, str):
        text = [(text, {})]
    for t, kw in text:
        r = p.add_run()
        r.text = t
        f = r.font
        f.size = Pt(kw.get("size", size))
        f.name = kw.get("font", font)
        f.bold = kw.get("bold", bold)
        f.italic = kw.get("italic", italic)
        f.color.rgb = kw.get("color", color)
    return p


def bullets(tf, items, size=16, color=INK, space_after=8, first=True,
            marker="▪", marker_color=ACCENT, line_spacing=1.05):
    for i, item in enumerate(items):
        lead = [(marker + "  ", {"color": marker_color, "size": size - 2, "bold": True})]
        if isinstance(item, str):
            body = [(item, {})]
        else:
            body = item
        put(tf, lead + body, size=size, color=color, space_after=space_after,
            first=(first and i == 0), line_spacing=line_spacing)


def chrome(slide, title, kicker=None):
    """Standard content-slide chrome: kicker, title, rule, page number."""
    global slide_no
    slide_no += 1
    if kicker:
        _, tf = textbox(slide, Inches(0.6), Inches(0.32), Inches(11), Inches(0.3))
        put(tf, kicker.upper(), size=12, color=ACCENT, bold=True, font=SANS, first=True)
    _, tf = textbox(slide, Inches(0.6), Inches(0.58), Inches(12.1), Inches(0.85))
    put(tf, title, size=30, color=INK, bold=True, font=SERIF, first=True)
    rect(slide, Inches(0.62), Inches(1.38), Inches(1.5), Pt(3), ACCENT)
    _, tf = textbox(slide, Inches(12.35), Inches(7.05), Inches(0.7), Inches(0.3))
    put(tf, str(slide_no), size=11, color=INK_SOFT, align=PP_ALIGN.RIGHT, first=True)
    _, tf = textbox(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3))
    put(tf, "VICTOR — a prompt-defined, cross-interface application framework",
        size=10, color=LINE if False else INK_SOFT, first=True)


def new_slide():
    return prs.slides.add_slide(BLANK)


def code_panel(slide, x, y, w, h, lines, size=12, title=None):
    rect(slide, x, y, w, h, CODE_BG, round_=True)
    ty = y + Inches(0.14)
    if title:
        _, tf = textbox(slide, x + Inches(0.25), ty, w - Inches(0.5), Inches(0.28))
        put(tf, title, size=11, color=CODE_HL, bold=True, font=MONO, first=True)
        ty += Inches(0.32)
    _, tf = textbox(slide, x + Inches(0.25), ty, w - Inches(0.5), h - (ty - y) - Inches(0.14))
    first = True
    for ln in lines:
        if isinstance(ln, str):
            ln = [(ln, {})]
        runs = [(t, {**{"font": MONO, "color": CODE_FG, "size": size}, **kw}) for t, kw in ln]
        put(tf, runs, size=size, color=CODE_FG, font=MONO, space_after=2,
            first=first, line_spacing=1.0)
        first = False


def chevron_flow(slide, x, y, w, h, steps, size=12.5, fill=ACCENT, gap=Inches(0.08)):
    n = len(steps)
    step_w = int((w - gap * (n - 1)) / n)
    cx = x
    for i, label in enumerate(steps):
        shape_kind = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        sp = slide.shapes.add_shape(shape_kind, cx, y, step_w, h)
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill if i % 2 == 0 else ACCENT2
        sp.line.fill.background()
        _shadow_off(sp)
        tf = sp.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = Inches(0.02)
        for j, line in enumerate(label.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(0)
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.name = SANS
            r.font.color.rgb = PAPER
        cx += step_w + gap


def card(slide, x, y, w, h, heading, body_items, head_color=ACCENT,
         head_size=15, body_size=12.5, fill=MIST):
    rect(slide, x, y, w, h, fill, round_=True)
    rect(slide, x, y, Inches(0.09), h, head_color, round_=False)
    _, tf = textbox(slide, x + Inches(0.28), y + Inches(0.16), w - Inches(0.5), Inches(0.4))
    put(tf, heading, size=head_size, color=INK, bold=True, font=SANS, first=True)
    _, tf = textbox(slide, x + Inches(0.28), y + Inches(0.58), w - Inches(0.5), h - Inches(0.75))
    first = True
    for item in body_items:
        if isinstance(item, str):
            item = [(item, {})]
        put(tf, item, size=body_size, color=INK_SOFT, space_after=5, first=first,
            line_spacing=1.02)
        first = False


# ================================================================== SLIDE 1 — title
s = new_slide()
slide_no += 1
rect(s, 0, 0, SLIDE_W, SLIDE_H, INK)
rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), ACCENT)
rect(s, Inches(0.9), Inches(2.62), Inches(2.2), Pt(3.5), AMBER)

_, tf = textbox(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.5))
put(tf, "COSMOPOLE  ·  ELPIAN PROJECT", size=15, color=CODE_HL, bold=True, first=True)

_, tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11.6), Inches(1.1))
put(tf, "VICTOR", size=66, color=PAPER, bold=True, font=SERIF, first=True)

_, tf = textbox(s, Inches(0.9), Inches(2.9), Inches(11.6), Inches(1.2))
put(tf, "A Prompt-Defined Programming Language and", size=24, color=PAPER,
    font=SERIF, first=True, space_after=2)
put(tf, "a Multi-Modal, Cross-Interface Application Framework", size=24,
    color=PAPER, font=SERIF, space_after=2)

_, tf = textbox(s, Inches(0.9), Inches(4.35), Inches(11.4), Inches(1.2))
put(tf, "From natural-language behaviour descriptions to platform-native executables:",
    size=16, color=CODE_FG, first=True, space_after=2)
put(tf, "modules · behaviours · views & actors · the Elpian bytecode VM",
    size=16, color=CODE_HL, space_after=2)

_, tf = textbox(s, Inches(0.9), Inches(5.9), Inches(11.4), Inches(0.8))
put(tf, "Cosmopole Research — github.com/cosmopole-org/victor", size=14,
    color=CODE_FG, first=True, space_after=2)
put(tf, "July 2026", size=14, color=CODE_FG, space_after=2)

# ================================================================== SLIDE 2 — outline
s = new_slide()
chrome(s, "Outline", "Agenda")
items_l = [
    [("1.  Motivation", {"bold": True, "color": INK}), ("  —  why fixed syntax and per-platform UIs no longer scale", {})],
    [("2.  The Victor programming model", {"bold": True, "color": INK}), ("  —  modules, behaviours, entrypoint", {})],
    [("3.  A language without fixed syntax", {"bold": True, "color": INK}), ("  —  human-language pseudo-code", {})],
    [("4.  Compilation pipeline", {"bold": True, "color": INK}), ("  —  prompts → operation families → AST → bytecode", {})],
    [("5.  The Elpian VM", {"bold": True, "color": INK}), ("  —  a no-JIT, governed execution substrate", {})],
]
items_r = [
    [("6.  Interface abstraction", {"bold": True, "color": INK}), ("  —  views and actors", {})],
    [("7.  Multi-modal builds", {"bold": True, "color": INK}), ("  —  CLI, voice, GUI, AR/VR from one definition", {})],
    [("8.  Implementation status", {"bold": True, "color": INK}), ("  —  what runs today", {})],
    [("9.  Case study", {"bold": True, "color": INK}), ("  —  VICTOR: CITY STRIKE (3D game on the VM)", {})],
    [("10. Related work, roadmap and conclusions", {"bold": True, "color": INK})],
]
_, tf = textbox(s, Inches(0.8), Inches(2.0), Inches(5.9), Inches(4.6))
bullets(tf, items_l, size=16, color=INK_SOFT, space_after=22)
_, tf = textbox(s, Inches(7.0), Inches(2.0), Inches(5.9), Inches(4.6))
bullets(tf, items_r, size=16, color=INK_SOFT, space_after=22)

# ================================================================== SLIDE 3 — motivation
s = new_slide()
chrome(s, "Motivation: three couplings that limit software", "1 · Motivation")
card(s, Inches(0.6), Inches(1.75), Inches(4.0), Inches(2.6),
     "Logic ⟷ Syntax", [
        "Programs must be expressed in a fixed grammar; the developer translates intent into syntax by hand.",
        "The abstraction ceiling is the language, not the problem.",
        [("Victor: ", {"bold": True, "color": ACCENT}), ("behaviour is written as structured natural language.", {})],
     ])
card(s, Inches(4.75), Inches(1.75), Inches(4.0), Inches(2.6),
     "Logic ⟷ Interface", [
        "UI code is written per modality: a GUI, a CLI, a voice skill and an AR scene are four codebases.",
        "Interaction assumptions leak into core logic.",
        [("Victor: ", {"bold": True, "color": ACCENT2}), ("interfaces are derived from abstract views and actors.", {})],
     ], head_color=ACCENT2)
card(s, Inches(8.9), Inches(1.75), Inches(4.0), Inches(2.6),
     "Logic ⟷ Platform", [
        "Each target (iOS, web, desktop, game engine) imposes its own toolchain and update rules (e.g. no JIT on iOS).",
        "Shipping one behaviour everywhere multiplies engineering cost.",
        [("Victor: ", {"bold": True, "color": AMBER}), ("one bytecode, executed by the Elpian VM on every target.", {})],
     ], head_color=AMBER)

rect(s, Inches(0.6), Inches(4.6), Inches(12.3), Inches(1.95), INK, round_=True)
_, tf = textbox(s, Inches(0.95), Inches(4.82), Inches(11.6), Inches(1.6))
put(tf, "Thesis", size=14, color=CODE_HL, bold=True, first=True, space_after=4)
put(tf, [("An application is fully specified by ", {}),
         ("what its modules do", {"bold": True, "color": PAPER}),
         (" (behaviours), ", {}),
         ("how they can be interacted with", {"bold": True, "color": PAPER}),
         (" (abstract inputs/outputs), and ", {}),
         ("where execution starts", {"bold": True, "color": PAPER}),
         (" (the entrypoint). Everything else — syntax, widgets, voice grammars, AR geometry, platform binaries — can be derived mechanically.", {})],
    size=16.5, color=CODE_FG, space_after=0, line_spacing=1.15)

# ================================================================== SLIDE 4 — programming model
s = new_slide()
chrome(s, "The Victor programming model", "2 · Programming model")
_, tf = textbox(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(0.5))
put(tf, [("A Victor program is a set of ", {}),
         ("modules", {"bold": True, "color": ACCENT}),
         (", each defined entirely by prompts. The developer decomposes the application, then describes each part.", {})],
    size=16, color=INK_SOFT, first=True)

card(s, Inches(0.6), Inches(2.25), Inches(3.95), Inches(3.15),
     "① Module decomposition", [
        "The application is split into cohesive modules (e.g. cart, inventory, checkout).",
        [("Internal modules", {"bold": True, "color": INK}), (" carry pure logic and state.", {})],
        [("Edge modules", {"bold": True, "color": INK}), (" own user-facing input/output and are interface-adapted per build.", {})],
     ])
card(s, Inches(4.7), Inches(2.25), Inches(3.95), Inches(3.15),
     "② Behaviour prompts", [
        "Each module carries a prompt that explains:",
        [("· behaviour", {"bold": True, "color": INK}), (" — what the module does, step by step;", {})],
        [("· contract", {"bold": True, "color": INK}), (" — the inputs it accepts and the outputs it produces for each behaviour.", {})],
        "Interconnections between modules are declared alongside the behaviours.",
     ], head_color=ACCENT2)
card(s, Inches(8.8), Inches(2.25), Inches(3.95), Inches(3.15),
     "③ Entrypoint", [
        "The developer designates the entry module/behaviour where execution begins.",
        "From there, the compiler resolves the module graph, type-checks the contracts and schedules the program.",
     ], head_color=AMBER)

rect(s, Inches(0.6), Inches(5.65), Inches(12.3), Inches(1.0), MIST, round_=True)
_, tf = textbox(s, Inches(0.95), Inches(5.86), Inches(11.7), Inches(0.7))
put(tf, [("Key property — consumer-aware outputs: ", {"bold": True, "color": ACCENT}),
         ("a module's output is not fixed at authoring time. The framework detects the consumer of each behaviour — a human on a screen, a voice channel, a terminal, an AR scene, or another module — and materialises the output in the form that consumer requires.", {})],
    size=14.5, color=INK_SOFT, first=True, line_spacing=1.1)

# ================================================================== SLIDE 5 — language
s = new_slide()
chrome(s, "A language without fixed syntax", "3 · Language design")
_, tf = textbox(s, Inches(0.6), Inches(1.6), Inches(5.7), Inches(0.9))
put(tf, [("Victor programs are ", {}),
         ("structured natural language", {"bold": True, "color": ACCENT}),
         (": numbered steps of human-readable pseudo-code. The wording is free; only the ", {}),
         ("intent", {"italic": True}), (" must be recoverable.", {})],
    size=15.5, color=INK_SOFT, first=True, line_spacing=1.1)

_, tf = textbox(s, Inches(0.6), Inches(2.85), Inches(5.7), Inches(3.8))
bullets(tf, [
    [("Dynamic, flexible surface form", {"bold": True, "color": INK}),
     (" — “define @a as 0”, “memorize 0 as @a” and “create a counter @a starting at zero” compile identically.", {})],
    [("Steps, not statements", {"bold": True, "color": INK}),
     (" — control flow references step numbers (“if …, go to step 5”), mirroring the stack-based instruction list underneath.", {})],
    [("Sub-steps", {"bold": True, "color": INK}),
     (" — parenthesised operations split into sub-steps; independent ones may be scheduled in parallel.", {})],
    [("Deterministic core", {"bold": True, "color": INK}),
     (" — every step is normalised into one of a closed set of operation families (next slide), so execution semantics stay exact.", {})],
], size=14.5, color=INK_SOFT, space_after=12)

code_panel(s, Inches(6.6), Inches(1.7), Inches(6.15), Inches(5.0), [
    [("program [", {"color": CODE_HL})],
    [("  step 1. ", {"color": AMBER}), ("function sum for inputs (operand1, operand2) [", {})],
    [("    step 1. ", {"color": AMBER}), ("define a new variable and name it @a", {})],
    [("            with initial value 0.", {})],
    [("    step 2. ", {"color": AMBER}), ("calculate operand1 + operand2", {})],
    [("            and assign it to @a.", {})],
    [("    step 3. ", {"color": AMBER}), ("return the value of @a.", {})],
    [("  ]", {})],
    [("  step 2. ", {"color": AMBER}), ("do the sum function on 2 and 3", {})],
    [("          and assign the result to @b.", {})],
    [("  step 3. ", {"color": AMBER}), ("call host to println variable @b.", {})],
    [("]", {"color": CODE_HL})],
    [("", {})],
    [("// same program, different words — same AST:", {"color": RGBColor(0x8A, 0x99, 0xB5)})],
    [("//   \"memorize zero as @a\"  ≡  \"define @a = 0\"", {"color": RGBColor(0x8A, 0x99, 0xB5)})],
], size=12.5, title="a Victor source program (pending-work/prompts/type_checker.prm)")

# ================================================================== SLIDE 6 — operation families
s = new_slide()
chrome(s, "Normalisation: nine operation families", "4 · Compilation")
_, tf = textbox(s, Inches(0.6), Inches(1.55), Inches(12.2), Inches(0.55))
put(tf, [("A classifier system-prompt assigns every step to one ", {}),
         ("operation family", {"bold": True, "color": ACCENT}),
         ("; a family-specific extractor prompt then populates that family's template AST object (JSON).", {})],
    size=15, color=INK_SOFT, first=True)

fams = [
    ("1  definition", "introduce new data in memory — “define @count with initial value 0”"),
    ("2  assignment", "update an existing variable — “update @x's value to y”"),
    ("3  calculation", "arithmetic / expression trees — “calculate (x + 1) ^ 2”"),
    ("4  functionDefinition", "declare a callable body of steps with parameters"),
    ("5  conditionalBranch", "test a condition; jump to true/false step numbers"),
    ("6  jumpOperation", "unconditional jump to a step in the same scope"),
    ("7  returnOperation", "return a value from the current function"),
    ("8  functionCall", "invoke a previously defined function"),
    ("9  host_call", "ask the host environment — println, getInput, initApp, …"),
]
x0, y0 = Inches(0.6), Inches(2.25)
cw, ch = Inches(4.03), Inches(1.22)
gx, gy = Inches(0.11), Inches(0.12)
for i, (name, desc) in enumerate(fams):
    r_, c_ = divmod(i, 3)
    x = x0 + c_ * (cw + gx)
    y = y0 + r_ * (ch + gy)
    rect(s, x, y, cw, ch, MIST, round_=True)
    _, tf = textbox(s, x + Inches(0.18), y + Inches(0.1), cw - Inches(0.36), Inches(0.32))
    put(tf, name, size=14, color=ACCENT, bold=True, font=MONO, first=True)
    _, tf = textbox(s, x + Inches(0.18), y + Inches(0.44), cw - Inches(0.36), ch - Inches(0.52))
    put(tf, desc, size=12, color=INK_SOFT, first=True, line_spacing=1.02)

rect(s, Inches(0.6), Inches(6.35), Inches(12.3), Inches(0.62), INK, round_=True)
_, tf = textbox(s, Inches(0.95), Inches(6.47), Inches(11.7), Inches(0.42))
put(tf, [("Values share one ", {}), ("Val", {"font": MONO, "color": CODE_HL, "bold": True}),
         (" structure: i16 / i32 / i64 · f32 / f64 · bool · string · object · array · function · identifier — matching the Elpian VM's tagged value model exactly.", {})],
    size=13.5, color=CODE_FG, first=True)

# ================================================================== SLIDE 7 — pipeline
s = new_slide()
chrome(s, "Compilation pipeline: prose to executable", "4 · Compilation")
chevron_flow(s, Inches(0.55), Inches(1.75), Inches(12.35), Inches(1.05), [
    "Prompts",
    "Classify",
    "AST\n(JSON)",
    "Program",
    "Bytecode",
    "Run",
], size=13)

steps_detail = [
    ("Classification", "Each step of every module prompt is classified by a system prompt into one of the nine operation families (assignment, arithmetic, function definition, …).", ACCENT),
    ("Template population", "A per-family extractor prompt pulls the operands, names and step targets out of the description and writes them into the prepared JSON AST template's fields.", ACCENT2),
    ("Program assembly", "Populated AST objects are pushed, in order, onto the program's stack-based instruction list; modules are linked through their declared interconnections and the entrypoint.", AMBER),
    ("Compile & run", "The final AST is handed to the Elpian compiler, lowered to bytecode, and executed by the Elpian VM inside the target environment (terminal, browser, mobile shell, 3D engine, …).", ACCENT),
]
x0, y0 = Inches(0.6), Inches(3.2)
cw = Inches(6.05); ch = Inches(1.55)
for i, (h, b, c) in enumerate(steps_detail):
    r_, c_ = divmod(i, 2)
    card(s, x0 + c_ * (cw + Inches(0.2)), y0 + r_ * (ch + Inches(0.18)), cw, ch,
         f"{i+1} · {h}", [b], head_color=c, body_size=12.5)

_, tf = textbox(s, Inches(0.6), Inches(6.62), Inches(12.2), Inches(0.4))
put(tf, [("The same pipeline runs once per ", {}), ("build type", {"bold": True, "color": ACCENT}),
         (" — each declaring the target environment's accessibility features (available input/output devices) — so one source produces many platform-specific executables.", {})],
    size=13.5, color=INK_SOFT, first=True)

# ================================================================== SLIDE 8 — Elpian VM
s = new_slide()
chrome(s, "The execution substrate: the Elpian VM", "5 · Elpian VM")
_, tf = textbox(s, Inches(0.6), Inches(1.6), Inches(6.2), Inches(0.6))
put(tf, [("A ", {}), ("pausing bytecode interpreter", {"bold": True, "color": ACCENT}),
         (" — no JIT, no machine-code generation — designed for places dynamic code is otherwise forbidden.", {})],
    size=15, color=INK_SOFT, first=True, line_spacing=1.1)

_, tf = textbox(s, Inches(0.6), Inches(2.4), Inches(6.2), Inches(4.3))
bullets(tf, [
    [("App-Store-legal & web-ready", {"bold": True, "color": INK}),
     (" — never writes executable memory (no W^X violation); compiles to wasm32 and native.", {})],
    [("One host seam: ", {"bold": True, "color": INK}), ("askHost(api, payload)", {"font": MONO, "color": ACCENT2}),
     (" — the VM suspends and hands every environment effect to the embedder; this is where Victor's interface adaptation plugs in.", {})],
    [("Governed by construction", {"bold": True, "color": INK}),
     (" — capability families (Gpu, Network, Storage, Clock, Randomness, …), instruction/memory/call-depth limits, resource meters that fail closed.", {})],
    [("Signed code delivery", {"bold": True, "color": INK}),
     (" — SHA-256/HMAC verify-before-load with downgrade protection: behaviour updates ship as data, safely.", {})],
    [("Multi-VM trees", {"bold": True, "color": INK}),
     (" — a VM can spawn sandboxed child VMs (module isolation); permissions AND down the ancestor path, budgets aggregate over subtrees.", {})],
], size=14, color=INK_SOFT, space_after=11)

code_panel(s, Inches(7.1), Inches(1.7), Inches(5.65), Inches(3.62), [
    [("[ Victor modules — bytecode ]", {"color": CODE_HL, "bold": True})],
    [("        |", {})],
    [("        v  runs on", {})],
    [("[ Elpian VM ]", {"color": CODE_HL, "bold": True})],
    [("   - stack machine", {})],
    [("   - tagged Vals (i16 ... f64)", {})],
    [("   - capability governor", {})],
    [("        |", {})],
    [("        v  askHost(api, payload)", {"color": AMBER})],
    [("[ Host bridges — per target ]", {"color": CODE_HL, "bold": True})],
    [("   terminal | Flutter/Skia", {})],
    [("   CanvasKit (web) | Godot 3D", {})],
], size=12.5, title="execution stack")

rect(s, Inches(7.1), Inches(5.48), Inches(5.65), Inches(1.32), MIST, round_=True)
_, tf = textbox(s, Inches(7.4), Inches(5.62), Inches(5.1), Inches(1.1))
put(tf, [("Front-ends already target it: ", {"bold": True, "color": INK})],
    size=13, first=True, space_after=3)
put(tf, [("js2elpian", {"font": MONO, "color": ACCENT}), (" and ", {}),
         ("dart2elpian", {"font": MONO, "color": ACCENT}),
         (" compile JS / a Dart subset to the same bytecode — Victor's prompt front-end is the third, syntax-free front-end.", {})],
    size=13, color=INK_SOFT, space_after=0, line_spacing=1.08)

# ================================================================== SLIDE 9 — view & actor
s = new_slide()
chrome(s, "Interface abstraction: views and actors", "6 · Views & actors")
_, tf = textbox(s, Inches(0.6), Inches(1.55), Inches(12.2), Inches(0.6))
put(tf, [("An edge module describes its interface as a ", {}),
         ("hierarchical XML/JSON tree", {"bold": True, "color": ACCENT}),
         (". Every node is one of two abstract types, each carrying a behaviour prompt — never a concrete widget.", {})],
    size=15, color=INK_SOFT, first=True)

card(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(2.2), "view — non-interactive output", [
    [("Shows data from the app's backbone logic or its data store, as the target modality allows:", {})],
    [("text · typography · image · chart · a spoken description · a spinning 3D cube · a JSON payload readable by another module", {"italic": True, "color": INK})],
], head_color=ACCENT2)
card(s, Inches(6.9), Inches(2.2), Inches(6.0), Inches(2.2), "actor — interactive input", [
    [("Accepts interactions from an interactor — a human or another module:", {})],
    [("button · checkbox · text field · a question awaiting a voice command · a key press · an AR gesture · a JSON message from a peer module", {"italic": True, "color": INK})],
], head_color=AMBER)

code_panel(s, Inches(0.6), Inches(4.6), Inches(12.3), Inches(2.35), [
    [("{ \"module\": \"checkout\",  \"interface\": {", {})],
    [("    \"type\": \"view\",  \"prompt\": \"show the running order total from the cart store\",", {"color": CODE_HL})],
    [("    \"children\": [", {})],
    [("      { \"type\": \"view\",  \"prompt\": \"list each cart item with its price\" },", {"color": CODE_HL})],
    [("      { \"type\": \"actor\", \"prompt\": \"let the interactor confirm the purchase\" },", {"color": RGBColor(0xFC, 0xD3, 0x4D)})],
    [("      { \"type\": \"actor\", \"prompt\": \"let the interactor change an item's quantity\" } ] } }", {"color": RGBColor(0xFC, 0xD3, 0x4D)})],
], size=12.5, title="abstract interface tree (authored once, adapted per modality)")

# ================================================================== SLIDE 10 — multimodal adaptation
s = new_slide()
chrome(s, "One definition, every modality", "7 · Multi-modal builds")
_, tf = textbox(s, Inches(0.6), Inches(1.55), Inches(12.2), Inches(0.55))
put(tf, [("Per build type, the compiler scans the ", {}),
         ("edge modules", {"bold": True, "color": ACCENT}),
         (" and adapts each view/actor to the declared I/O devices of the target environment:", {})],
    size=15, color=INK_SOFT, first=True)

# adaptation matrix
rows = [
    ("Terminal / CLI", "formatted text, tables, ASCII charts", "flags, prompts, key presses", ACCENT),
    ("Voice / vocal systems", "spoken descriptions, earcons", "voice commands, spoken answers", ACCENT2),
    ("GUI (desktop & mobile)", "typography, images, charts, dashboards", "buttons, checkboxes, text fields", AMBER),
    ("AR / VR floating interfaces", "3D objects, spatial panels, animated meshes", "gestures, gaze, controller input", ACCENT),
    ("3D / game environments", "in-world HUDs, objects (e.g. a spinning cube)", "in-world interactions, touch controls", ACCENT2),
    ("Module-to-module", "machine-readable JSON payloads", "JSON messages received from peers", AMBER),
]
x0, y0 = Inches(0.6), Inches(2.3)
col_w = [Inches(3.3), Inches(4.7), Inches(4.3)]
hh = Inches(0.42)
hdrs = ["Target environment", "views become …", "actors become …"]
cx = x0
for j, h in enumerate(hdrs):
    rect(s, cx, y0, col_w[j], hh, INK)
    _, tf = textbox(s, cx + Inches(0.15), y0 + Inches(0.07), col_w[j] - Inches(0.3), hh - Inches(0.1))
    put(tf, h, size=13, color=PAPER, bold=True, first=True)
    cx += col_w[j]
rh = Inches(0.6)
for i, (env, v, a, c) in enumerate(rows):
    y = y0 + hh + i * rh
    bg = MIST if i % 2 == 0 else PAPER
    cx = x0
    for j, txt in enumerate([env, v, a]):
        rect(s, cx, y, col_w[j], rh, bg, line_color=LINE, line_w=Pt(0.5))
        _, tf = textbox(s, cx + Inches(0.15), y + Inches(0.08), col_w[j] - Inches(0.3), rh - Inches(0.12))
        put(tf, txt, size=12.5, color=(INK if j == 0 else INK_SOFT),
            bold=(j == 0), first=True, line_spacing=0.95)
        cx += col_w[j]

_, tf = textbox(s, Inches(0.6), Inches(6.55), Inches(12.2), Inches(0.5))
put(tf, [("The interactor need not be human: ", {"bold": True, "color": ACCENT}),
         ("when a behaviour's consumer is another module, the same view/actor pair degrades to a typed JSON contract — UI and IPC are one abstraction.", {})],
    size=13.5, color=INK_SOFT, first=True)

# ================================================================== SLIDE 11 — architecture overview
s = new_slide()
chrome(s, "End-to-end architecture", "7 · Multi-modal builds")

# left column: authoring
rect(s, Inches(0.6), Inches(1.7), Inches(3.6), Inches(4.9), MIST, round_=True)
_, tf = textbox(s, Inches(0.85), Inches(1.85), Inches(3.1), Inches(0.35))
put(tf, "AUTHORING", size=13, color=ACCENT, bold=True, first=True)
for i, t in enumerate(["module A — behaviour prompt", "module B — behaviour prompt",
                       "edge module — view/actor tree", "entrypoint declaration"]):
    y = Inches(2.3) + i * Inches(1.02)
    rect(s, Inches(0.85), y, Inches(3.1), Inches(0.84), PAPER, line_color=LINE, round_=True)
    _, tf = textbox(s, Inches(1.0), y + Inches(0.12), Inches(2.85), Inches(0.66))
    put(tf, t, size=12.5, color=INK, first=True, line_spacing=1.0)

# middle column: compiler
rect(s, Inches(4.75), Inches(1.7), Inches(3.6), Inches(4.9), INK, round_=True)
_, tf = textbox(s, Inches(5.0), Inches(1.85), Inches(3.1), Inches(0.35))
put(tf, "VICTOR COMPILER", size=13, color=CODE_HL, bold=True, first=True)
for i, t in enumerate(["classify → 9 operation families", "populate template AST (JSON)",
                       "assemble instruction list + link modules", "Elpian compiler → bytecode"]):
    y = Inches(2.3) + i * Inches(1.02)
    rect(s, Inches(5.0), y, Inches(3.1), Inches(0.84), RGBColor(0x24, 0x33, 0x4E), round_=True)
    _, tf = textbox(s, Inches(5.15), y + Inches(0.12), Inches(2.85), Inches(0.66))
    put(tf, t, size=12.5, color=CODE_FG, first=True, line_spacing=1.0)

# right column: targets
rect(s, Inches(8.9), Inches(1.7), Inches(3.85), Inches(4.9), MIST, round_=True)
_, tf = textbox(s, Inches(9.15), Inches(1.85), Inches(3.4), Inches(0.35))
put(tf, "BUILD TYPES → TARGETS", size=13, color=AMBER, bold=True, first=True)
tgt = [("CLI build", "terminal I/O"), ("Voice build", "speech in/out"),
       ("GUI build", "Flutter/Skia · CanvasKit web"), ("AR/VR & 3D build", "Godot engine bridge")]
for i, (a, b) in enumerate(tgt):
    y = Inches(2.3) + i * Inches(1.02)
    rect(s, Inches(9.15), y, Inches(3.35), Inches(0.84), PAPER, line_color=LINE, round_=True)
    _, tf = textbox(s, Inches(9.3), y + Inches(0.1), Inches(3.1), Inches(0.7))
    put(tf, a, size=12.5, color=INK, bold=True, first=True, space_after=1)
    put(tf, b, size=11, color=INK_SOFT, space_after=0)

# arrows
for xa in (Inches(4.28), Inches(8.43)):
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, xa, Inches(3.85), Inches(0.42), Inches(0.5))
    _fill(ar, ACCENT)
    _shadow_off(ar)

_, tf = textbox(s, Inches(0.6), Inches(6.72), Inches(12.2), Inches(0.35))
put(tf, [("Every target runs the ", {}), ("same bytecode on the same VM", {"bold": True, "color": ACCENT}),
         (" — only the host bridge behind askHost differs.", {})],
    size=13.5, color=INK_SOFT, first=True)

# ================================================================== SLIDE 12 — governance
s = new_slide()
chrome(s, "Governance: execution remains provable", "5 · Elpian VM")
_, tf = textbox(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(0.55))
put(tf, [("Natural-language authoring raises an obvious concern — ", {}),
         ("what guarantees the produced program is safe?", {"italic": True, "color": INK}),
         ("  Victor answers at the execution layer, not the language layer:", {})],
    size=15, color=INK_SOFT, first=True)

card(s, Inches(0.6), Inches(2.35), Inches(6.0), Inches(1.95), "Two-layer capability model", [
    [("VM backstop:", {"bold": True, "color": INK}), (" coarse families — Gpu, Network, Storage, Clock, Randomness — a disabled family short-circuits to a typed null.", {})],
    [("Host layer:", {"bold": True, "color": INK}), (" per-library capabilities that fail closed for unknown surfaces; sandboxed profiles deny I/O by default.", {})],
], head_color=ACCENT)
card(s, Inches(6.9), Inches(2.35), Inches(6.0), Inches(1.95), "Resource metering", [
    "Instruction, memory and call-depth limits; host-call counts and bytes-across-the-seam are bounded per module.",
    "In multi-VM trees, budgets are enforced against aggregate subtree usage — a runaway module cannot starve its siblings.",
], head_color=ACCENT2)
card(s, Inches(0.6), Inches(4.5), Inches(6.0), Inches(1.95), "Verified delivery", [
    "Bundles are SHA-256/HMAC signed, verified before load, with downgrade protection and content-hash-pinned manifests.",
    "Behaviour updates are data — deployable live even where JIT is forbidden (iOS, web).",
], head_color=AMBER)
card(s, Inches(6.9), Inches(4.5), Inches(6.0), Inches(1.95), "Deterministic semantics", [
    "Free-form prose is normalised to a closed instruction set before execution — the VM never interprets natural language at runtime.",
    "The prompt layer is a compiler front-end; auditing the emitted AST/bytecode is always possible.",
], head_color=ACCENT)

# ================================================================== SLIDE 13 — implementation status
s = new_slide()
chrome(s, "Implementation status in this repository", "8 · Status")
rows = [
    ("Elpian VM (victor/elpian-vm)", "AST/bytecode pausing interpreter; capability + resource governor; native & wasm32", "✅ built & tested"),
    ("Language front-ends (js2elpian, dart2elpian)", "JS and Dart-subset compilers to Elpian bytecode; universal stdlib names resolved at compile time", "✅ built & tested"),
    ("Dart runtime layer (victor/dart)", "dart:ui, typed_data, convert, isolates, async loop, widget layer, flutter.dart library", "✅ 190 tests pass"),
    ("Web renderer (victor/web-demo)", "reflective CanvasKit/Skia bridge — 575 symbols audited, 0 unreachable; headless-browser verified", "✅ built & verified"),
    ("Godot bridge (victor/bridge)", "reflective controller over ClassDB (~900 classes); multi-VM tree with sandboxing & budgets", "✅ built & verified"),
    ("Prompt → AST pipeline (pending-work/prompts)", "classifier + 9 per-family extractor prompts emitting template AST JSON", "🔶 prototype"),
    ("Interface adaptation (views/actors per modality)", "edge-module scanning and per-device materialisation", "🔷 design stage"),
]
x0, y0 = Inches(0.6), Inches(1.8)
col_w = [Inches(4.35), Inches(5.9), Inches(2.05)]
hh = Inches(0.42)
cx = x0
for j, h in enumerate(["Component", "What it does", "Status"]):
    rect(s, cx, y0, col_w[j], hh, INK)
    _, tf = textbox(s, cx + Inches(0.13), y0 + Inches(0.07), col_w[j] - Inches(0.26), hh - Inches(0.1))
    put(tf, h, size=13, color=PAPER, bold=True, first=True)
    cx += col_w[j]
rh = Inches(0.64)
for i, (a, b, c) in enumerate(rows):
    y = y0 + hh + i * rh
    bg = MIST if i % 2 == 0 else PAPER
    cx = x0
    for j, txt in enumerate([a, b, c]):
        rect(s, cx, y, col_w[j], rh, bg, line_color=LINE, line_w=Pt(0.5))
        _, tf = textbox(s, cx + Inches(0.13), y + Inches(0.06), col_w[j] - Inches(0.26), rh - Inches(0.1))
        put(tf, txt, size=11.5, color=(INK if j == 0 else INK_SOFT), bold=(j == 0),
            first=True, line_spacing=0.95)
        cx += col_w[j]

_, tf = textbox(s, Inches(0.6), Inches(6.72), Inches(12.2), Inches(0.35))
put(tf, "Honest scope: the substrate and bridges are real and test-covered; the prompt front-end and adaptation are the research edge.",
    size=13, color=INK_SOFT, italic=True, first=True)

# ================================================================== SLIDE 14 — case study
s = new_slide()
chrome(s, "Case study — VICTOR: CITY STRIKE", "9 · Case study")
_, tf = textbox(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(0.6))
put(tf, [("A complete third-person shooter written entirely in Dart and executed on the Elpian VM inside Godot 4 — evidence the substrate can carry ", {}),
         ("full interactive 3D applications", {"bold": True, "color": ACCENT}),
         (", the richest target modality Victor compiles to.", {})],
    size=15, color=INK_SOFT, first=True, line_spacing=1.1)

_, tf = textbox(s, Inches(0.6), Inches(2.45), Inches(6.1), Inches(4.2))
bullets(tf, [
    [("Procedurally assembled city", {"bold": True, "color": INK}), (" from CC0 GLB asset kits.", {})],
    [("Animated player character", {"bold": True, "color": INK}), (" with over-the-shoulder SpringArm camera.", {})],
    [("Two hitscan weapons", {"bold": True, "color": INK}), (" with pooled tracers, impacts and damage numbers.", {})],
    [("Three enemy archetypes", {"bold": True, "color": INK}), (" — chase / attack / line-of-sight AI, wave spawning, pickups.", {})],
    [("Full HUD and menus, touch controls, synthesized embedded audio.", {})],
    [("Verified headless end-to-end", {"bold": True, "color": INK}), (" by the capi test suite (run_tps.rs).", {})],
], size=14, color=INK_SOFT, space_after=10)

rect(s, Inches(7.1), Inches(2.45), Inches(5.65), Inches(4.2), INK, round_=True)
_, tf = textbox(s, Inches(7.4), Inches(2.65), Inches(5.1), Inches(0.4))
put(tf, "Why it matters for Victor", size=15, color=CODE_HL, bold=True, first=True)
_, tf = textbox(s, Inches(7.4), Inches(3.1), Inches(5.1), Inches(3.4))
bullets(tf, [
    [("The Godot bridge is reflective", {"bold": True, "color": PAPER}),
     (" — it addresses every engine class/method/property by name, so coverage is complete by construction.", {})],
    [("The multi-VM tree", {"bold": True, "color": PAPER}),
     (" maps directly onto Victor modules: one sandboxed VM per module, with metered budgets and AND-composed permissions.", {})],
    [("The same program shape", {"bold": True, "color": PAPER}),
     (" (bytecode + askHost) will render a dashboard on web, a HUD in 3D, or text in a terminal.", {})],
], size=13, color=CODE_FG, space_after=10, marker_color=AMBER)

# ================================================================== SLIDE 15 — related work
s = new_slide()
chrome(s, "Position among related work", "10 · Related work")
rows = [
    ("Cross-platform UI frameworks", "Flutter, React Native, Compose MP",
     "Share logic and widgets across screens, but stay within one modality (GUI) and one fixed-syntax language.",
     "Victor shares behaviour across modalities; GUI is just one build type."),
    ("LLM code assistants", "Copilot, Cursor, agentic coding",
     "Generate conventional source code that developers then own, in a fixed-syntax language.",
     "In Victor the prompt is the source; generation targets a closed AST, not free-form code."),
    ("Model-driven engineering", "UML/DSL-to-code generators",
     "Abstract models compile to implementations, but models are rigid, diagram- or grammar-bound.",
     "Victor's 'model' is flexible natural language, normalised by classification."),
    ("Intent-based / declarative UI", "SwiftUI, HTMX, voice-skill DSLs",
     "Declare what the interface shows per platform toolkit.",
     "Victor declares interface meaning (view/actor) independent of any toolkit."),
]
x0, y0 = Inches(0.6), Inches(1.75)
col_w = [Inches(2.7), Inches(2.3), Inches(3.9), Inches(3.4)]
hh = Inches(0.42)
cx = x0
for j, h in enumerate(["Field", "Examples", "What it gives", "How Victor differs"]):
    rect(s, cx, y0, col_w[j], hh, INK)
    _, tf = textbox(s, cx + Inches(0.12), y0 + Inches(0.07), col_w[j] - Inches(0.24), hh - Inches(0.1))
    put(tf, h, size=12.5, color=PAPER, bold=True, first=True)
    cx += col_w[j]
rh = Inches(1.08)
for i, row in enumerate(rows):
    y = y0 + hh + i * rh
    bg = MIST if i % 2 == 0 else PAPER
    cx = x0
    for j, txt in enumerate(row):
        rect(s, cx, y, col_w[j], rh, bg, line_color=LINE, line_w=Pt(0.5))
        _, tf = textbox(s, cx + Inches(0.12), y + Inches(0.08), col_w[j] - Inches(0.24), rh - Inches(0.14))
        put(tf, txt, size=11, color=(INK if j == 0 else (ACCENT if j == 3 else INK_SOFT)),
            bold=(j == 0), first=True, line_spacing=0.98)
        cx += col_w[j]

_, tf = textbox(s, Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.35))
put(tf, [("Victor's claim: ", {"bold": True, "color": ACCENT}),
         ("behaviour, interface and platform are independent axes — only behaviour is authored by hand.", {})],
    size=13.5, color=INK_SOFT, first=True)

# ================================================================== SLIDE 16 — roadmap
s = new_slide()
chrome(s, "Roadmap", "10 · Roadmap")
phases = [
    ("Now", "Solidify the prompt front-end", [
        "Harden the classifier + extractor prompts (pending-work/prompts) into a reproducible compiler stage with regression corpora.",
        "Emit Elpian AST directly from populated templates; round-trip audits of prose → AST → prose.",
    ], ACCENT),
    ("Next", "Interface adaptation engine", [
        "Formalise the view/actor tree schema (XML/JSON) and the per-build device-capability descriptors.",
        "Materialisers: terminal renderer, flutter.dart GUI, CanvasKit web, Godot 3D/AR — all through the existing askHost bridges.",
    ], ACCENT2),
    ("Later", "Full multi-modal releases", [
        "Voice and AR/VR materialisers; consumer detection for module-to-module JSON contracts.",
        "One-command builds: every declared build type → signed, governed, platform-specific executable.",
    ], AMBER),
]
x0 = Inches(0.6)
cw = Inches(4.03)
for i, (tag, head, items, c) in enumerate(phases):
    x = x0 + i * (cw + Inches(0.12))
    rect(s, x, Inches(1.85), cw, Inches(0.5), c, round_=True)
    _, tf = textbox(s, x, Inches(1.93), cw, Inches(0.35))
    put(tf, tag.upper(), size=14, color=PAPER, bold=True, align=PP_ALIGN.CENTER, first=True)
    rect(s, x, Inches(2.5), cw, Inches(3.9), MIST, round_=True)
    _, tf = textbox(s, x + Inches(0.25), Inches(2.7), cw - Inches(0.5), Inches(0.6))
    put(tf, head, size=15.5, color=INK, bold=True, first=True, line_spacing=1.0)
    _, tf = textbox(s, x + Inches(0.25), Inches(3.35), cw - Inches(0.5), Inches(2.9))
    bullets(tf, items, size=12.5, color=INK_SOFT, space_after=10, marker_color=c)

_, tf = textbox(s, Inches(0.6), Inches(6.68), Inches(12.2), Inches(0.35))
put(tf, "Each stage lands as a tested vertical slice — the method that shipped Elpian phases 1–8.",
    size=13, color=INK_SOFT, italic=True, first=True)

# ================================================================== SLIDE 17 — conclusions
s = new_slide()
chrome(s, "Conclusions", "Conclusions")
_, tf = textbox(s, Inches(0.6), Inches(1.75), Inches(12.2), Inches(3.4))
bullets(tf, [
    [("Programs as behaviour, not syntax. ", {"bold": True, "color": INK}),
     ("Victor programs are module prompts — flexible human language normalised into nine exact operation families, then compiled like any language.", {})],
    [("Interfaces as meaning, not widgets. ", {"bold": True, "color": INK}),
     ("Two abstract element types — view and actor — describe every interface; per-build adaptation materialises them for terminals, voice, GUIs, AR/VR, 3D worlds, or peer modules.", {})],
    [("Execution as a governed substrate. ", {"bold": True, "color": INK}),
     ("The Elpian VM makes the whole scheme deployable — no JIT, capability-governed, resource-metered, signature-verified — on iOS, web, desktop, mobile and game engines alike.", {})],
    [("The substrate is proven. ", {"bold": True, "color": INK}),
     ("190 passing tests, a Flutter-style widget stack rasterised by real Skia, and a complete 3D shooter running on the VM inside Godot.", {})],
], size=16, color=INK_SOFT, space_after=16, line_spacing=1.1)

rect(s, Inches(0.6), Inches(5.35), Inches(12.3), Inches(1.35), INK, round_=True)
_, tf = textbox(s, Inches(0.95), Inches(5.62), Inches(11.6), Inches(0.9))
put(tf, [("Write what each module should do. Declare how it can be talked to. Pick an entrypoint.", {"bold": True, "color": PAPER, "size": 18, "font": SERIF})],
    size=18, first=True, space_after=4)
put(tf, [("Victor derives the rest — for every interface the user owns.", {"color": CODE_HL, "size": 15, "italic": True})], size=15, space_after=0)

# ================================================================== SLIDE 18 — references / thanks
s = new_slide()
slide_no += 1
rect(s, 0, 0, SLIDE_W, SLIDE_H, INK)
rect(s, Inches(0.9), Inches(1.95), Inches(2.2), Pt(3.5), AMBER)
_, tf = textbox(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.9))
put(tf, "Thank you", size=46, color=PAPER, bold=True, font=SERIF, first=True)

_, tf = textbox(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(0.5))
put(tf, "References & artefacts", size=16, color=CODE_HL, bold=True, first=True)

refs = [
    ("Repository", "github.com/cosmopole-org/victor — this presentation: presentation/"),
    ("Elpian VM & Dart layer", "victor/ — VM, dart2elpian, js2elpian, flutter.dart, 190-test suite (victor/README.md)"),
    ("Elpian VM upstream", "github.com/cosmopole-org/elpis (crates/elpian-vm)"),
    ("Prompt → AST pipeline", "pending-work/prompts/ — classifier + nine operation-family extractor prompts"),
    ("Web / Skia bridge", "victor/web-demo/canvaskit_bridge.js — reflective CanvasKit driver (575 symbols audited)"),
    ("Godot bridge & case study", "victor/bridge/ — reflective ClassDB controller, multi-VM tree, VICTOR: CITY STRIKE (GAME_DESIGN.md)"),
]
_, tf = textbox(s, Inches(0.9), Inches(2.9), Inches(11.6), Inches(3.4))
for i, (a, b) in enumerate(refs):
    put(tf, [(a + "  —  ", {"bold": True, "color": PAPER}), (b, {"color": CODE_FG})],
        size=14, first=(i == 0), space_after=12)

_, tf = textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5))
put(tf, "Cosmopole Research · July 2026", size=13, color=CODE_FG, first=True)

# ------------------------------------------------------------------ save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Victor-App-Framework.pptx")
prs.save(out)
print(f"wrote {out} ({slide_no} slides)")
